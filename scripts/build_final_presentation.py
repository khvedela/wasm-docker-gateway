#!/usr/bin/env python3
"""Generate final_presentation.pptx from benchmark outputs and a PPT template."""

from __future__ import annotations

import argparse
import csv
import datetime as dt
import os
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


EXPECTED_VARIANTS = [
    "native_local",
    "native_docker",
    "wasm_host_cli",
    "wasm_host_wasmtime",
    "wasm_host_wasmtime_embedded",
]

VARIANT_LABELS = {
    "native_local": "Native (local)",
    "native_docker": "Native (Docker)",
    "wasm_host_cli": "Wasm host (CLI spawn)",
    "wasm_host_wasmtime": "Wasm host (Wasmtime CLI)",
    "wasm_host_wasmtime_embedded": "Wasm host (Wasmtime embedded)",
}

VARIANT_SHORT = {
    "native_local": "local",
    "native_docker": "docker",
    "wasm_host_cli": "cli-spawn",
    "wasm_host_wasmtime": "wasmtime-cli",
    "wasm_host_wasmtime_embedded": "embedded",
}


@dataclass
class DeckData:
    cold_by_variant: Dict[str, Dict[str, str]] = field(default_factory=dict)
    warm_by_variant: Dict[str, Dict[str, str]] = field(default_factory=dict)
    throughput_rows: List[Dict[str, str]] = field(default_factory=list)
    variants_seen: List[str] = field(default_factory=list)
    env_snapshot_path: Optional[Path] = None
    env_snapshot_lines: List[str] = field(default_factory=list)
    missing_inputs: List[str] = field(default_factory=list)
    fallback_notes: List[str] = field(default_factory=list)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate final benchmark presentation deck.")
    parser.add_argument(
        "--template",
        default="",
        help="Path to .pptx/.potx template. If omitted, auto-detect from repo root or templates/.",
    )
    parser.add_argument("--results-dir", default="results", help="Results root directory.")
    parser.add_argument("--output", default="final_presentation.pptx", help="Output .pptx filename/path.")
    parser.add_argument("--student", default="[Student name]", help="Student name.")
    parser.add_argument("--tutors", default="[Tutor names]", help="Tutor names.")
    parser.add_argument("--date", default="", help="Presentation date (example: Feb 23 2026).")
    return parser.parse_args()


def read_csv_rows(path: Path) -> List[Dict[str, str]]:
    if not path.exists():
        return []
    with path.open("r", encoding="utf-8", newline="") as f:
        return list(csv.DictReader(f))


def as_float(value: str, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def fmt_num(value: float, digits: int = 2) -> str:
    if value == 0:
        return "0"
    if digits == 0:
        return str(int(round(value)))
    return f"{value:.{digits}f}".rstrip("0").rstrip(".")


def fmt_rps(value: float) -> str:
    return fmt_num(value, 0 if value >= 100 else 1)


def safe_ratio(numerator: float, denominator: float) -> Optional[float]:
    if denominator <= 0:
        return None
    return numerator / denominator


def kb_to_mb(kb: float) -> float:
    return kb / 1024.0


def short_text(text: str, max_chars: int) -> str:
    text = text.strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3] + "..."


def normalize_date(date_arg: str) -> str:
    if not date_arg.strip():
        return dt.date.today().strftime("%B %d, %Y")
    cleaned = date_arg.strip()
    for fmt in ("%b %d %Y", "%B %d %Y", "%Y-%m-%d"):
        try:
            parsed = dt.datetime.strptime(cleaned, fmt)
            return parsed.strftime("%B %d, %Y")
        except ValueError:
            continue
    return cleaned


def choose_template(workspace: Path, explicit: str) -> Path:
    if explicit:
        path = Path(explicit).expanduser().resolve()
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {path}")
        return path

    candidates: List[Path] = []
    for pattern in ("*.pptx", "*.potx", "templates/*.pptx", "templates/*.potx"):
        candidates.extend((workspace / pattern).parent.glob(Path(pattern).name))

    desktop_fallback = Path("/Users/david/Desktop/template-cedric.potx")
    if desktop_fallback.exists():
        candidates.append(desktop_fallback)

    unique: Dict[str, Path] = {}
    for c in candidates:
        unique[str(c.resolve())] = c.resolve()
    deduped = list(unique.values())

    if not deduped:
        raise FileNotFoundError("No template found (.pptx/.potx).")

    def score(path: Path) -> Tuple[int, int]:
        has_template_name = 1 if "template" in path.name.lower() else 0
        return (has_template_name, int(path.stat().st_size))

    return sorted(deduped, key=score, reverse=True)[0]


def convert_potx_for_pptx(template_path: Path) -> Tuple[Path, Optional[Path]]:
    """python-pptx cannot open .potx directly; convert package content-type."""
    if template_path.suffix.lower() != ".potx":
        return template_path, None

    fd, temp_name = tempfile.mkstemp(prefix="template_", suffix=".pptx")
    os.close(fd)
    tmp_path = Path(temp_name)

    template_ct = (
        b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
    )
    presentation_ct = (
        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
    )

    with zipfile.ZipFile(template_path, "r") as zin, zipfile.ZipFile(
        tmp_path, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for info in zin.infolist():
            payload = zin.read(info.filename)
            if info.filename == "[Content_Types].xml":
                payload = payload.replace(template_ct, presentation_ct)
            zout.writestr(info, payload)

    return tmp_path, tmp_path


def clear_all_slides(prs: Presentation) -> None:
    for sld_id in list(prs.slides._sldIdLst):  # type: ignore[attr-defined]
        rel_id = sld_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(sld_id)  # type: ignore[attr-defined]


def force_text_size(text_frame, size_pt: int) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""
        for run in paragraph.runs:
            if run.text and run.text.strip():
                run.font.size = Pt(size_pt)


def set_title(slide, title: str, subtitle: Optional[str] = None) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        force_text_size(slide.shapes.title.text_frame, 34)
    if subtitle is not None and len(slide.placeholders) > 1:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                placeholder.text = subtitle
                force_text_size(placeholder.text_frame, 20)
                break


def set_bullets(shape, lines: Sequence[str], size_pt: int = 18) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
    force_text_size(text_frame, size_pt)


def add_textbox(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    size_pt: int = 17,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    if not paragraph.runs:
        paragraph.add_run()
    for run in paragraph.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def set_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    notes.notes_text_frame.clear()
    notes.notes_text_frame.text = text


def add_picture_fit(slide, image_path: Path, left: int, top: int, width: int, height: int):
    picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
    if picture.height > height:
        scale = height / picture.height
        picture.height = int(picture.height * scale)
        picture.width = int(picture.width * scale)
    picture.left = left + int((width - picture.width) / 2)
    picture.top = top + int((height - picture.height) / 2)
    return picture


def pick_plot(
    plots_dir: Path,
    expected_name: str,
    fallbacks: Sequence[str],
    fallback_notes: List[str],
) -> Optional[Path]:
    expected = plots_dir / expected_name
    if expected.exists():
        return expected

    for candidate in fallbacks:
        p = plots_dir / candidate
        if p.exists():
            fallback_notes.append(f"Missing {expected_name}; used {candidate}.")
            return p

    all_png = sorted(plots_dir.glob("*.png"))
    if all_png:
        fallback_notes.append(f"Missing {expected_name}; used {all_png[0].name}.")
        return all_png[0]

    fallback_notes.append(f"Missing {expected_name}; no fallback available.")
    return None


def find_env_snapshot(results_dir: Path) -> Optional[Path]:
    exact = results_dir / "meta" / "env_snapshot.txt"
    if exact.exists():
        return exact
    meta_dir = results_dir / "meta"
    if not meta_dir.exists():
        return None
    candidates = sorted(meta_dir.glob("env_snapshot*.txt"))
    return candidates[-1] if candidates else None


def parse_env_snapshot_lines(path: Path) -> List[str]:
    raw_lines = [line.strip() for line in path.read_text(encoding="utf-8", errors="replace").splitlines()]
    collected: List[str] = []
    mem_line = ""
    for line in raw_lines:
        if not line:
            continue
        lower = line.lower()
        if lower.startswith("timestamp:"):
            collected.append(line)
        elif lower.startswith("kernel:"):
            collected.append(short_text(line, 85))
        elif lower.startswith("cpu_cores:"):
            collected.append(line)
        elif lower.startswith("git_commit:"):
            collected.append(short_text(line, 65))
        elif lower.startswith("rustc:"):
            collected.append(short_text(line, 65))
        elif line.startswith("Mem:") and not mem_line:
            parts = line.split()
            if len(parts) >= 8:
                mem_line = f"memory: total {parts[1]}, available {parts[7]}"
            else:
                mem_line = short_text(line, 65)

    if mem_line:
        collected.append(mem_line)

    # Keep concise and readable in slide body.
    return collected[:5]


def live_system_snapshot_lines() -> List[str]:
    lines: List[str] = []
    timestamp = dt.datetime.now().isoformat(timespec="seconds")
    lines.append(f"timestamp: {timestamp}")
    kernel = subprocess.getoutput("uname -a").strip()
    if kernel:
        lines.append(short_text(f"kernel: {kernel}", 85))
    cores = subprocess.getoutput("sysctl -n hw.ncpu 2>/dev/null").strip()
    if not cores:
        cores = str(os.cpu_count() or "")
    if cores:
        lines.append(f"cpu_cores: {cores}")
    mem_bytes = subprocess.getoutput("sysctl -n hw.memsize 2>/dev/null").strip()
    if mem_bytes.isdigit():
        gib = int(mem_bytes) / (1024**3)
        lines.append(f"memory: total {gib:.1f} GiB")
    return lines[:5]


def load_data(results_dir: Path) -> DeckData:
    data = DeckData()
    summary_dir = results_dir / "summary"

    cold_rows = read_csv_rows(summary_dir / "cold_start_percentiles.csv")
    warm_rows = read_csv_rows(summary_dir / "warm_latency_percentiles.csv")
    throughput_rows = read_csv_rows(summary_dir / "throughput_summary.csv")

    if not cold_rows:
        data.missing_inputs.append("results/summary/cold_start_percentiles.csv")
    if not warm_rows:
        data.missing_inputs.append("results/summary/warm_latency_percentiles.csv")
    if not throughput_rows:
        data.missing_inputs.append("results/summary/throughput_summary.csv")

    data.cold_by_variant = {row.get("variant", ""): row for row in cold_rows if row.get("variant")}
    data.warm_by_variant = {row.get("variant", ""): row for row in warm_rows if row.get("variant")}
    data.throughput_rows = throughput_rows

    variants_found = set(data.cold_by_variant.keys()) | set(data.warm_by_variant.keys())
    variants_found |= {row.get("variant", "") for row in throughput_rows if row.get("variant")}
    data.variants_seen = [v for v in EXPECTED_VARIANTS if v in variants_found]

    env_path = find_env_snapshot(results_dir)
    data.env_snapshot_path = env_path
    if env_path:
        data.env_snapshot_lines = parse_env_snapshot_lines(env_path)
    else:
        data.env_snapshot_lines = live_system_snapshot_lines()
        data.missing_inputs.append("results/meta/env_snapshot.txt")

    return data


def select_throughput_row(
    throughput_rows: Sequence[Dict[str, str]],
    workload: str,
    variant: str,
    target_conns: int,
) -> Optional[Dict[str, str]]:
    candidates = [
        row
        for row in throughput_rows
        if row.get("workload") == workload and row.get("variant") == variant
    ]
    if not candidates:
        return None

    exact = [row for row in candidates if int(as_float(row.get("conns", "0"))) == target_conns]
    if exact:
        return exact[0]

    return min(
        candidates,
        key=lambda row: abs(int(as_float(row.get("conns", "0"))) - target_conns),
    )


def select_best_throughput_row(
    throughput_rows: Sequence[Dict[str, str]],
    workload: str,
    variant: str,
) -> Optional[Dict[str, str]]:
    candidates = [
        row
        for row in throughput_rows
        if row.get("workload") == workload and row.get("variant") == variant
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda row: as_float(row.get("mean_rps", "0")))


def sorted_metric_rows(by_variant: Dict[str, Dict[str, str]], metric: str) -> List[Tuple[str, Dict[str, str]]]:
    items = list(by_variant.items())
    return sorted(items, key=lambda kv: as_float(kv[1].get(metric, ""), float("inf")))


def variant_label(variant: str) -> str:
    return VARIANT_LABELS.get(variant, variant)


def variant_short(variant: str) -> str:
    return VARIANT_SHORT.get(variant, variant)


def cold_takeaways(data: DeckData) -> List[str]:
    rows = sorted_metric_rows(data.cold_by_variant, "p50_ms")
    if len(rows) < 2:
        return ["Cold start summary is unavailable for at least two variants."]

    fast_variant, fast_row = rows[0]
    slow_variant, slow_row = rows[-1]
    fast_p50 = as_float(fast_row.get("p50_ms", "0"))
    slow_p50 = as_float(slow_row.get("p50_ms", "0"))
    ratio = safe_ratio(slow_p50, fast_p50)

    line1 = (
        f"Fastest p50 cold start: {variant_short(fast_variant)} "
        f"{fmt_num(fast_p50, 1)} ms."
    )
    if ratio is None:
        line2 = f"Slowest p50: {variant_short(slow_variant)} {fmt_num(slow_p50, 1)} ms."
    else:
        line2 = (
            f"Slowest p50: {variant_short(slow_variant)} {fmt_num(slow_p50, 1)} ms "
            f"({fmt_num(ratio, 1)}x slower than {variant_short(fast_variant)})."
        )

    wasm_entries = []
    for v in ("wasm_host_wasmtime", "wasm_host_cli", "wasm_host_wasmtime_embedded"):
        row = data.cold_by_variant.get(v)
        if row:
            wasm_entries.append(f"{variant_short(v)} {fmt_num(as_float(row.get('p50_ms', '0')), 1)}")
    line3 = ""
    if len(wasm_entries) == 3:
        line3 = "Wasm cold p50 ordering: " + " < ".join(wasm_entries) + " ms."

    return [line for line in (line1, line2, line3) if line]


def warm_takeaways(data: DeckData) -> List[str]:
    rows = sorted_metric_rows(data.warm_by_variant, "p50_ms")
    if len(rows) < 2:
        return ["Warm latency summary is unavailable for at least two variants."]

    base_variant, base_row = rows[0]
    base_p50 = as_float(base_row.get("p50_ms", "0"))

    native = data.warm_by_variant.get("native_local")
    docker = data.warm_by_variant.get("native_docker")
    embedded = data.warm_by_variant.get("wasm_host_wasmtime_embedded")
    wt_cli = data.warm_by_variant.get("wasm_host_wasmtime")
    cli_spawn = data.warm_by_variant.get("wasm_host_cli")

    line1 = (
        f"Top warm p50 group: local {fmt_num(as_float(native.get('p50_ms', '0')) if native else 0, 2)} ms, "
        f"docker {fmt_num(as_float(docker.get('p50_ms', '0')) if docker else 0, 2)} ms, "
        f"embedded {fmt_num(as_float(embedded.get('p50_ms', '0')) if embedded else 0, 2)} ms."
    )

    wt_p50 = as_float(wt_cli.get("p50_ms", "0")) if wt_cli else 0
    cli_p50 = as_float(cli_spawn.get("p50_ms", "0")) if cli_spawn else 0
    wt_factor = safe_ratio(wt_p50, base_p50) or 0
    cli_factor = safe_ratio(cli_p50, base_p50) or 0
    line2 = (
        f"CLI modes add fixed overhead: wasmtime-cli {fmt_num(wt_factor, 2)}x and "
        f"cli-spawn {fmt_num(cli_factor, 2)}x vs {variant_short(base_variant)} p50."
    )
    return [line1, line2]


def throughput_takeaways(data: DeckData, target_conns: int = 50) -> List[str]:
    proxy_parts: List[str] = []
    compute_parts: List[str] = []
    for variant in EXPECTED_VARIANTS:
        proxy_row = select_throughput_row(data.throughput_rows, "proxy", variant, target_conns)
        compute_row = select_throughput_row(data.throughput_rows, "compute", variant, target_conns)
        if proxy_row:
            proxy_parts.append(
                f"{variant_short(variant)} {fmt_rps(as_float(proxy_row.get('mean_rps', '0')))}"
            )
        if compute_row:
            compute_parts.append(
                f"{variant_short(variant)} {fmt_rps(as_float(compute_row.get('mean_rps', '0')))}"
            )

    line1 = "Proxy @50 conns RPS: " + ", ".join(proxy_parts) + "."
    line2 = "Compute @50 conns RPS: " + ", ".join(compute_parts) + "."

    local_proxy = select_throughput_row(data.throughput_rows, "proxy", "native_local", target_conns)
    embedded_proxy = select_throughput_row(
        data.throughput_rows, "proxy", "wasm_host_wasmtime_embedded", target_conns
    )
    cli_proxy = select_throughput_row(data.throughput_rows, "proxy", "wasm_host_cli", target_conns)
    if local_proxy and embedded_proxy and cli_proxy:
        local_rps = as_float(local_proxy.get("mean_rps", "0"))
        embedded_rps = as_float(embedded_proxy.get("mean_rps", "0"))
        cli_rps = as_float(cli_proxy.get("mean_rps", "0"))
        embedded_share = safe_ratio(embedded_rps, local_rps)
        cli_share = safe_ratio(cli_rps, local_rps)
        if embedded_share and cli_share:
            line3 = (
                f"Proxy throughput share vs local: embedded {fmt_num(embedded_share * 100, 1)}%, "
                f"cli-spawn {fmt_num(cli_share * 100, 1)}%."
            )
            return [line1, line2, line3]

    return [line1, line2]


def efficiency_takeaways(data: DeckData, target_conns: int = 50) -> List[str]:
    cpu_parts: List[Tuple[str, float]] = []
    rss_parts: List[str] = []

    for variant in EXPECTED_VARIANTS:
        row = select_throughput_row(data.throughput_rows, "proxy", variant, target_conns)
        if not row:
            continue
        cpu_eff = as_float(row.get("cpu_per_1k_rps", "0"))
        rss_mb = kb_to_mb(as_float(row.get("mean_rss_kb", "0")))
        cpu_parts.append((variant_short(variant), cpu_eff))
        rss_parts.append(f"{variant_short(variant)} {fmt_num(rss_mb, 1)} MB")

    cpu_parts.sort(key=lambda x: x[1])
    cpu_text = ", ".join([f"{name} {fmt_num(value, 1)}" for name, value in cpu_parts])
    line1 = f"CPU per 1k RPS (proxy @50, lower is better): {cpu_text}."
    line2 = f"RSS (proxy @50): {', '.join(rss_parts)}."
    return [line1, line2]


def add_architecture_diagram(slide) -> None:
    y = Inches(1.95)
    h = Inches(1.12)
    w = Inches(2.35)
    gap = Inches(0.42)
    x = Inches(0.7)

    labels = ["Load generator", "Gateway", "Handler", "Upstream"]
    boxes = []
    for label in labels:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.text_frame.clear()
        box.text_frame.text = label
        force_text_size(box.text_frame, 18)
        boxes.append(box)
        x += w + gap

    for left_box, right_box in zip(boxes, boxes[1:]):
        x1 = left_box.left + left_box.width
        y1 = left_box.top + int(left_box.height / 2)
        x2 = right_box.left
        y2 = right_box.top + int(right_box.height / 2)
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        connector.line.width = Pt(2)

    add_textbox(
        slide,
        "Workloads: hello, compute, state, proxy",
        Inches(0.8),
        Inches(3.72),
        Inches(6.5),
        Inches(0.55),
        size_pt=18,
        bold=True,
    )
    add_textbox(
        slide,
        "Metrics: cold start, warm latency (p50/p90/p99), throughput vs conns, RSS, CPU efficiency",
        Inches(0.8),
        Inches(4.35),
        Inches(12.0),
        Inches(0.75),
        size_pt=16,
    )


def build_deck(
    prs: Presentation,
    template_path: Path,
    output_path: Path,
    results_dir: Path,
    student: str,
    tutors: str,
    presentation_date: str,
) -> DeckData:
    data = load_data(results_dir)
    plots_dir = results_dir / "plots"
    if not plots_dir.exists():
        raise FileNotFoundError(f"Missing plots directory: {plots_dir}")

    cold_plot = pick_plot(
        plots_dir,
        "cold_start_median_p90.png",
        ["warm_latency_median_p90.png", "latency_compute.png"],
        data.fallback_notes,
    )
    warm_plot = pick_plot(
        plots_dir,
        "warm_latency_median_p90.png",
        ["latency_compute.png", "latency_proxy.png"],
        data.fallback_notes,
    )
    throughput_proxy_plot = pick_plot(
        plots_dir,
        "throughput_proxy.png",
        ["throughput_compute.png", "throughput_state.png"],
        data.fallback_notes,
    )
    throughput_compute_plot = pick_plot(
        plots_dir,
        "throughput_compute.png",
        ["throughput_proxy.png", "throughput_state.png"],
        data.fallback_notes,
    )
    efficiency_plot = pick_plot(
        plots_dir,
        "efficiency_proxy.png",
        ["rss_proxy.png", "efficiency_compute.png"],
        data.fallback_notes,
    )

    clear_all_slides(prs)
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_two = prs.slide_layouts[3]
    layout_title_only = prs.slide_layouts[5]

    # 1) Title
    slide = prs.slides.add_slide(layout_title)
    set_title(
        slide,
        "WebAssembly Runtime Comparison for Gateway Workloads",
        (
            "Native vs Docker vs Wasm runtimes\n"
            f"Tutors: {tutors} | Student: {student}\n"
            f"Date: {presentation_date}"
        ),
    )
    set_notes(
        slide,
        "Deck generated from benchmark artifacts.\n"
        f"Template: {template_path}\n"
        f"Results root: {results_dir}",
    )

    # 2) Problem/Motivation
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Problem & Motivation")
    set_bullets(
        slide.placeholders[1],
        [
            "Gateway handlers are short-lived and frequent, so fixed overhead can dominate tail latency.",
            "Cold start directly impacts burst handling and scale-to-zero behavior.",
            "Runtime choice affects latency, throughput, memory footprint, and operational complexity.",
            "Goal: provide data-backed guidance for native, Docker, and Wasm runtime models.",
        ],
        size_pt=18,
    )

    # 3) Platforms compared
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Platforms Compared (What Is Actually Measured)")
    set_bullets(
        slide.placeholders[1],
        [
            "native_local: host baseline for minimum overhead.",
            "native_docker: same gateway binary inside container runtime.",
            "wasm_host_cli: per-request runtime CLI process spawn.",
            "wasm_host_wasmtime: per-request Wasmtime CLI process path.",
            "wasm_host_wasmtime_embedded: in-process Wasmtime runtime (no spawn per request).",
        ],
        size_pt=17,
    )

    # 4) Architecture
    slide = prs.slides.add_slide(layout_title_only)
    set_title(slide, "Testbed Architecture & Workloads")
    add_architecture_diagram(slide)

    # 5) Methodology
    slide = prs.slides.add_slide(layout_two)
    set_title(slide, "Methodology / Experimental Protocol")
    conns = sorted(
        {int(as_float(row.get("conns", "0"))) for row in data.throughput_rows if row.get("conns")},
    )
    conns_text = ",".join(str(c) for c in conns) if conns else "n/a"
    cold_samples = (
        int(as_float(next(iter(data.cold_by_variant.values())).get("sample_count", "0")))
        if data.cold_by_variant
        else 0
    )
    warm_local = data.warm_by_variant.get("native_local")
    warm_samples = int(as_float(warm_local.get("sample_count", "0"))) if warm_local else 0

    set_bullets(
        slide.placeholders[1],
        [
            f"Cold start: hyperfine, warmup=0, sample_count={cold_samples} per variant.",
            f"Warm latency: hyperfine + curl, sample_count={warm_samples} (native_local).",
            f"Throughput grid: wrk over conns [{conns_text}], 30s windows.",
            "Resource tracking: sampler.py every 0.2s for RSS and CPU.",
        ],
        size_pt=16,
    )

    env_source = data.env_snapshot_path.name if data.env_snapshot_path else "live host fallback"
    env_lines = [f"Environment snapshot: {env_source}"] + data.env_snapshot_lines[:4]
    set_bullets(slide.placeholders[2], env_lines, size_pt=16)
    set_notes(
        slide,
        "Methodology source files: scripts/bench_cold_start.sh, scripts/bench_warm_latency.sh, "
        "scripts/bench_throughput.sh, scripts/bench_all.sh",
    )

    # 6) Cold start results
    slide = prs.slides.add_slide(layout_title_only)
    set_title(slide, "Cold Start Results")
    if cold_plot:
        add_picture_fit(slide, cold_plot, Inches(0.8), Inches(1.3), Inches(11.8), Inches(4.55))
        add_textbox(
            slide,
            f"Figure: {cold_plot.name}",
            Inches(0.85),
            Inches(5.9),
            Inches(6.0),
            Inches(0.4),
            size_pt=15,
        )
    add_textbox(
        slide,
        "\n".join(cold_takeaways(data)),
        Inches(0.8),
        Inches(6.1),
        Inches(11.8),
        Inches(0.85),
        size_pt=16,
    )
    set_notes(
        slide,
        "Cold summary rows:\n"
        + "\n".join(
            f"{variant}: p50={fmt_num(as_float(row.get('p50_ms', '0')), 2)} "
            f"p90={fmt_num(as_float(row.get('p90_ms', '0')), 2)} "
            f"p99={fmt_num(as_float(row.get('p99_ms', '0')), 2)}"
            for variant, row in sorted_metric_rows(data.cold_by_variant, "p50_ms")
        ),
    )

    # 7) Warm latency results
    slide = prs.slides.add_slide(layout_title_only)
    set_title(slide, "Warm Latency Results (p50 / p90)")
    if warm_plot:
        add_picture_fit(slide, warm_plot, Inches(0.8), Inches(1.3), Inches(11.8), Inches(4.55))
        add_textbox(
            slide,
            f"Figure: {warm_plot.name}",
            Inches(0.85),
            Inches(5.9),
            Inches(6.0),
            Inches(0.4),
            size_pt=15,
        )
    add_textbox(
        slide,
        "\n".join(warm_takeaways(data)),
        Inches(0.8),
        Inches(6.1),
        Inches(11.8),
        Inches(0.85),
        size_pt=16,
    )

    # 8) Throughput results
    slide = prs.slides.add_slide(layout_title_only)
    set_title(slide, "Throughput Results (RPS vs Connections)")
    if throughput_proxy_plot:
        add_picture_fit(slide, throughput_proxy_plot, Inches(0.6), Inches(1.35), Inches(6.1), Inches(3.85))
        add_textbox(slide, "Proxy workload", Inches(0.65), Inches(5.12), Inches(2.8), Inches(0.35), size_pt=15)
    if throughput_compute_plot:
        add_picture_fit(slide, throughput_compute_plot, Inches(6.75), Inches(1.35), Inches(6.1), Inches(3.85))
        add_textbox(slide, "Compute workload", Inches(6.8), Inches(5.12), Inches(3.3), Inches(0.35), size_pt=15)
    add_textbox(
        slide,
        "\n".join(throughput_takeaways(data, target_conns=50)),
        Inches(0.7),
        Inches(5.45),
        Inches(12.0),
        Inches(1.2),
        size_pt=16,
    )

    # 9) Efficiency/resources
    slide = prs.slides.add_slide(layout_title_only)
    set_title(slide, "Efficiency & Resource Footprint")
    if efficiency_plot:
        add_picture_fit(slide, efficiency_plot, Inches(0.8), Inches(1.35), Inches(7.65), Inches(4.9))
        add_textbox(slide, f"Figure: {efficiency_plot.name}", Inches(0.85), Inches(6.2), Inches(4.5), Inches(0.35), size_pt=15)
    add_textbox(
        slide,
        "\n".join(efficiency_takeaways(data, target_conns=50)),
        Inches(8.55),
        Inches(1.75),
        Inches(4.0),
        Inches(4.4),
        size_pt=16,
    )

    # 10) Key findings
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Key Findings")
    local_cold = as_float(data.cold_by_variant.get("native_local", {}).get("p50_ms", "0"))
    docker_cold = as_float(data.cold_by_variant.get("native_docker", {}).get("p50_ms", "0"))
    docker_cold_factor = safe_ratio(docker_cold, local_cold) or 0
    set_bullets(
        slide.placeholders[1],
        [
            f"Docker warm latency is close to local, but cold start is {fmt_num(docker_cold_factor, 1)}x slower in this testbed.",
            "Embedded Wasmtime is close to native on warm latency and compute throughput.",
            "CLI-spawn models show major warm-path penalties; they primarily measure process spawn overhead.",
            "Workload type changes rankings and gaps; there is no single best runtime for every constraint.",
        ],
        size_pt=17,
    )

    # 11) Recommendations
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Recommendations: When To Use What")
    set_bullets(
        slide.placeholders[1],
        [
            "Use native_local for strict latency SLOs and peak throughput goals.",
            "Use native_docker for portability/isolation when cold-start budget is not critical.",
            "Use Wasmtime embedded for Wasm portability in the hot path, with memory budget planning.",
            "Avoid per-request CLI runtime spawn in production gateways.",
        ],
        size_pt=18,
    )

    # 12) Threats / limitations
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Threats to Validity / Limitations")
    set_bullets(
        slide.placeholders[1],
        [
            "Single-host loopback setup may differ from real distributed network behavior.",
            "wrk is synthetic traffic; client limits and queueing can bias observed tails.",
            "Throughput summary has sample_count=1 per grid point in this run.",
            "Latency-under-load percentiles are not included in the current throughput summary.",
            "JIT/cache and host background load can shift absolute numbers between runs.",
        ],
        size_pt=16,
    )

    # 13) Summary / conclusion
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "Summary / Conclusion")
    proxy_local = select_throughput_row(data.throughput_rows, "proxy", "native_local", 50)
    proxy_embedded = select_throughput_row(
        data.throughput_rows, "proxy", "wasm_host_wasmtime_embedded", 50
    )
    summary_line = "Proxy @50 conns: "
    if proxy_local and proxy_embedded:
        summary_line += (
            f"local {fmt_rps(as_float(proxy_local.get('mean_rps', '0')))} rps, "
            f"embedded {fmt_rps(as_float(proxy_embedded.get('mean_rps', '0')))} rps."
        )
    else:
        summary_line += "insufficient data."

    set_bullets(
        slide.placeholders[1],
        [
            "Native local remains the raw performance baseline in this benchmark campaign.",
            "Embedded Wasmtime is the strongest Wasm option for warm-path gateway workloads.",
            "CLI-based Wasm execution is useful as a stress baseline for spawn overhead, not as a production data-plane target.",
            summary_line,
        ],
        size_pt=17,
    )

    # Optional appendix only when needed.
    appendix_lines = [f"Missing input: {item}" for item in sorted(set(data.missing_inputs))]
    appendix_lines.extend([f"Fallback used: {item}" for item in data.fallback_notes])
    if appendix_lines:
        slide = prs.slides.add_slide(layout_content)
        set_title(slide, "Appendix: Data Gaps / Fallbacks")
        set_bullets(slide.placeholders[1], appendix_lines[:8], size_pt=16)

    prs.save(str(output_path))
    return data


def print_readme_snippet(
    template_path: Path,
    output_path: Path,
    results_dir: Path,
    data: DeckData,
    student: str,
    tutors: str,
    date_text: str,
) -> None:
    print("\nREADME snippet")
    print("--------------")
    print("Rebuild command:")
    print(
        ".venv/bin/python scripts/build_final_presentation.py "
        f"--template \"{template_path}\" "
        f"--results-dir \"{results_dir}\" "
        f"--output \"{output_path.name}\" "
        f"--student \"{student}\" "
        f"--tutors \"{tutors}\" "
        f"--date \"{date_text}\""
    )
    print("Inputs used:")
    print(f"- template: {template_path}")
    print(f"- summaries: {results_dir / 'summary'}/*.csv")
    print(f"- plots: {results_dir / 'plots'}/*.png")
    env_ref = data.env_snapshot_path if data.env_snapshot_path else results_dir / "meta" / "env_snapshot.txt"
    print(f"- env snapshot reference: {env_ref}")


def main() -> None:
    args = parse_args()
    workspace = Path.cwd()
    results_dir = (workspace / args.results_dir).resolve()
    output_path = (workspace / args.output).resolve()
    date_text = normalize_date(args.date)

    template_path = choose_template(workspace, args.template)
    loadable_template, temp_created = convert_potx_for_pptx(template_path)

    try:
        prs = Presentation(str(loadable_template))
        data = build_deck(
            prs=prs,
            template_path=template_path,
            output_path=output_path,
            results_dir=results_dir,
            student=args.student,
            tutors=args.tutors,
            presentation_date=date_text,
        )
        print(f"Generated: {output_path}")
        print_readme_snippet(
            template_path=template_path,
            output_path=output_path,
            results_dir=results_dir,
            data=data,
            student=args.student,
            tutors=args.tutors,
            date_text=date_text,
        )
    finally:
        if temp_created and temp_created.exists():
            temp_created.unlink(missing_ok=True)


if __name__ == "__main__":
    main()
